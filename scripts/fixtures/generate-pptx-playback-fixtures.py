#!/usr/bin/env python3

"""生成体积小、结构明确的 PPTX 播放验收素材。"""

import argparse
from pathlib import Path
from tempfile import TemporaryDirectory
import wave
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from lxml import etree
import av
import numpy
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "tests/fixtures/pptx/playback-controlled.pptx"
NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p159": "http://schemas.microsoft.com/office/powerpoint/2015/09/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def shape_id(document: etree._Element, name: str) -> str:
    matches = document.xpath(f".//p:cNvPr[@name='{name}']", namespaces=NS)
    if len(matches) != 1:
        raise RuntimeError(f"对象名称不唯一：{name}")
    return matches[0].get("id")


def qname(prefix: str, local_name: str) -> str:
    return f"{{{NS[prefix]}}}{local_name}"


def transition(kind: str, duration_ms: int = 600, direction: str | None = None, advance_ms: int | None = None) -> etree._Element:
    attributes = {"dur": str(duration_ms), "advClick": "1"}
    if advance_ms is not None:
        attributes["advTm"] = str(advance_ms)
    if kind == "morph":
        alternate = etree.Element(qname("mc", "AlternateContent"), nsmap={
            "mc": NS["mc"],
            "p159": NS["p159"],
            "p14": NS["p14"],
        })
        choice = etree.SubElement(alternate, qname("mc", "Choice"), {"Requires": "p159"})
        choice_transition = etree.SubElement(choice, qname("p", "transition"), {"advClick": "1"})
        choice_transition.set(qname("p14", "dur"), str(duration_ms))
        etree.SubElement(choice_transition, qname("p159", "morph"), {"option": "byObject"})
        fallback = etree.SubElement(alternate, qname("mc", "Fallback"))
        fallback_transition = etree.SubElement(fallback, qname("p", "transition"), {"advClick": "1"})
        etree.SubElement(fallback_transition, qname("p", "fade"))
        return alternate
    node = etree.Element(qname("p", "transition"), attributes)
    effect = etree.SubElement(node, qname("p", kind))
    if direction:
        effect.set("dir", direction)
    return node


def timing(target_id: str, trigger_id: str) -> etree._Element:
    xml = f"""
    <p:timing xmlns:p="{NS['p']}">
      <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst><p:seq><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
          <p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
            <p:par><p:cTn id="4" presetID="10" presetClass="entr" nodeType="afterEffect" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
                <p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="5" dur="700"/>
                  <p:tgtEl><p:spTgt spid="{target_id}"/></p:tgtEl>
                </p:cBhvr></p:animEffect>
              </p:childTnLst>
            </p:cTn></p:par>
          </p:childTnLst></p:cTn></p:par>
          <p:par><p:cTn id="6"><p:stCondLst><p:cond evt="onClick" delay="0">
            <p:tgtEl><p:spTgt spid="{trigger_id}"/></p:tgtEl>
          </p:cond></p:stCondLst><p:childTnLst>
            <p:par><p:cTn id="7" presetID="10" presetClass="entr" nodeType="withEffect" fill="hold">
              <p:childTnLst><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="8" dur="500"/>
                <p:tgtEl><p:spTgt spid="{trigger_id}"/></p:tgtEl>
              </p:cBhvr></p:animEffect></p:childTnLst>
            </p:cTn></p:par>
          </p:childTnLst></p:cTn></p:par>
        </p:childTnLst></p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
        </p:seq></p:childTnLst>
      </p:cTn></p:par></p:tnLst>
    </p:timing>
    """
    return etree.fromstring(xml.encode())


def media_timing(video_id: str, audio_id: str, bookmark_target_id: str) -> etree._Element:
    xml = f"""
    <p:timing xmlns:p="{NS['p']}" xmlns:p14="{NS['p14']}">
      <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        <p:par><p:cTn id="2" presetClass="mediacall" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:cmd type="call" cmd="play"><p:cBhvr><p:cTn id="3" dur="1"/>
            <p:tgtEl><p:spTgt spid="{video_id}"/></p:tgtEl>
          </p:cBhvr></p:cmd>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="4"><p:stCondLst><p:cond evt="onMediaBookmark" delay="0">
          <p:tgtEl><p14:bmkTgt spid="{video_id}" bmkName="middle"/></p:tgtEl>
        </p:cond></p:stCondLst><p:childTnLst>
          <p:par><p:cTn id="5" presetID="10" presetClass="entr" nodeType="withEffect" fill="hold"><p:childTnLst>
            <p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="6" dur="300"/>
              <p:tgtEl><p:spTgt spid="{bookmark_target_id}"/></p:tgtEl>
            </p:cBhvr></p:animEffect>
          </p:childTnLst></p:cTn></p:par>
        </p:childTnLst></p:cTn></p:par>
        <p:video><p:cMediaNode vol="70000"><p:cTn id="7" dur="indefinite" repeatCount="indefinite"/>
          <p:tgtEl><p:spTgt spid="{video_id}"/></p:tgtEl>
        </p:cMediaNode></p:video>
        <p:par><p:cTn id="8" presetClass="mediacall" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:cmd type="call" cmd="play"><p:cBhvr><p:cTn id="9" dur="1"/>
            <p:tgtEl><p:spTgt spid="{audio_id}"/></p:tgtEl>
          </p:cBhvr></p:cmd>
        </p:childTnLst></p:cTn></p:par>
        <p:audio><p:cMediaNode vol="40000"><p:cTn id="10" dur="indefinite"/>
          <p:tgtEl><p:spTgt spid="{audio_id}"/></p:tgtEl>
        </p:cMediaNode></p:audio>
      </p:childTnLst></p:cTn></p:par></p:tnLst>
    </p:timing>
    """
    return etree.fromstring(xml.encode())


def repeat_timing(scale_id: str, motion_id: str, opacity_id: str) -> etree._Element:
    xml = f"""
    <p:timing xmlns:p="{NS['p']}">
      <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        <p:par><p:cTn id="2" presetClass="emph" nodeType="withEffect" repeatCount="2000" autoRev="1" fill="hold"><p:childTnLst>
          <p:animScale><p:cBhvr><p:cTn id="3" dur="500"/><p:tgtEl><p:spTgt spid="{scale_id}"/></p:tgtEl></p:cBhvr>
            <p:by x="120000" y="120000"/>
          </p:animScale>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="4" presetClass="path" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:animMotion path="M 0 0 L 0.25 -0.1 E"><p:cBhvr><p:cTn id="5" dur="600"/>
            <p:tgtEl><p:spTgt spid="{motion_id}"/></p:tgtEl>
          </p:cBhvr></p:animMotion>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="6" presetClass="emph" nodeType="withEffect" autoRev="1" fill="hold"><p:childTnLst>
          <p:anim><p:cBhvr><p:cTn id="7" dur="800"/><p:tgtEl><p:spTgt spid="{opacity_id}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
          </p:cBhvr><p:tavLst>
            <p:tav tm="0"><p:val><p:fltVal val="100000"/></p:val></p:tav>
            <p:tav tm="100000"><p:val><p:fltVal val="30000"/></p:val></p:tav>
          </p:tavLst></p:anim>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="8" presetClass="emph" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:animRot by="2700000"><p:cBhvr><p:cTn id="9" dur="600"/>
            <p:tgtEl><p:spTgt spid="{motion_id}"/></p:tgtEl>
          </p:cBhvr></p:animRot>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="10" presetID="22" presetClass="entr" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:animEffect transition="in" filter="wipe(right)"><p:cBhvr><p:cTn id="11" dur="600"/>
            <p:tgtEl><p:spTgt spid="{opacity_id}"/></p:tgtEl>
          </p:cBhvr></p:animEffect>
        </p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn></p:par></p:tnLst>
    </p:timing>
    """
    return etree.fromstring(xml.encode())


def indefinite_timing(scale_id: str) -> etree._Element:
    xml = f"""
    <p:timing xmlns:p="{NS['p']}"><p:tnLst><p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst><p:par>
        <p:cTn id="2" presetClass="emph" nodeType="withEffect" repeatCount="indefinite" autoRev="1" fill="hold"><p:childTnLst>
          <p:animScale><p:cBhvr><p:cTn id="3" dur="300"/><p:tgtEl><p:spTgt spid="{scale_id}"/></p:tgtEl></p:cBhvr>
            <p:by x="130000" y="130000"/>
          </p:animScale>
        </p:childTnLst></p:cTn>
      </p:par></p:childTnLst></p:cTn>
    </p:par></p:tnLst></p:timing>
    """
    return etree.fromstring(xml.encode())


def overlap_timing(target_id: str) -> etree._Element:
    xml = f"""
    <p:timing xmlns:p="{NS['p']}"><p:tnLst><p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        <p:par><p:cTn id="2" presetID="10" presetClass="entr" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="3" dur="500"/>
            <p:tgtEl><p:spTgt spid="{target_id}"/></p:tgtEl>
          </p:cBhvr></p:animEffect>
        </p:childTnLst></p:cTn></p:par>
        <p:par><p:cTn id="4" presetID="10" presetClass="exit" nodeType="withEffect" fill="hold"><p:childTnLst>
          <p:animEffect transition="out" filter="fade"><p:cBhvr><p:cTn id="5" dur="500"/>
            <p:tgtEl><p:spTgt spid="{target_id}"/></p:tgtEl>
          </p:cBhvr></p:animEffect>
        </p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn>
    </p:par></p:tnLst></p:timing>
    """
    return etree.fromstring(xml.encode())


def paragraph_timing(target_id: str) -> etree._Element:
    effects = []
    for paragraph in range(3):
        base = 2 + paragraph * 3
        effects.append(f"""
        <p:par><p:cTn id="{base}" nodeType="clickEffect"><p:childTnLst><p:par>
          <p:cTn id="{base + 1}" presetID="10" presetClass="entr" nodeType="withEffect" fill="hold"><p:childTnLst>
            <p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{base + 2}" dur="300"/>
              <p:tgtEl><p:spTgt spid="{target_id}"><p:txEl><p:pRg st="{paragraph}" end="{paragraph}"/></p:txEl></p:spTgt></p:tgtEl>
            </p:cBhvr></p:animEffect>
          </p:childTnLst></p:cTn>
        </p:par></p:childTnLst></p:cTn></p:par>
        """)
    xml = f"""
    <p:timing xmlns:p="{NS['p']}"><p:tnLst><p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        {''.join(effects)}
      </p:childTnLst></p:cTn>
    </p:par></p:tnLst></p:timing>
    """
    return etree.fromstring(xml.encode())


def add_slide(prs: Presentation, title: str, color: RGBColor, hero_name: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)
    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(3.2), Inches(2.1))
    hero.name = hero_name or f"hero-{title}"
    hero.fill.solid()
    hero.fill.fore_color.rgb = color
    hero.line.fill.background()
    text = slide.shapes.add_textbox(Inches(1), Inches(0.45), Inches(8), Inches(0.7))
    text.name = f"title-{title}"
    paragraph = text.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    trigger = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.8), Inches(1), Inches(1))
    trigger.name = f"trigger-{title}"
    trigger.fill.solid()
    trigger.fill.fore_color.rgb = RGBColor(15, 23, 42)


def named_shape(prs: Presentation, slide_index: int, name: str):
    return next(shape for shape in prs.slides[slide_index].shapes if shape.name == name)


def patch_archive(source: Path, target: Path, patched_slides: set[int]) -> None:
    replacements: dict[str, bytes] = {}
    with ZipFile(source) as archive:
        for index in range(1, 13):
            path = f"ppt/slides/slide{index}.xml"
            document = etree.fromstring(archive.read(path))
            if index not in patched_slides:
                continue
            # python-pptx 会为视频自动写入一棵媒体时间树。受控素材需要用包含
            # 播放命令和书签的完整时间树替换它；同一页不能保留两个 p:timing。
            for existing_timing in document.findall(qname("p", "timing")):
                document.remove(existing_timing)
            color_map = document.find(qname("p", "clrMapOvr"))
            insert_at = list(document).index(color_map) + 1 if color_map is not None else 1
            if index == 1:
                target_id = shape_id(document, "hero-单击动画")
                trigger_id = shape_id(document, "trigger-单击动画")
                document.insert(insert_at, timing(target_id, trigger_id))
            elif index == 2:
                document.insert(insert_at, transition("fade", 650))
            elif index == 3:
                document.insert(insert_at, transition("push", 700, "l"))
            elif index == 4:
                document.insert(insert_at, transition("wipe", 750, "r", 900))
            elif index == 5:
                document.insert(insert_at, transition("morph", 800))
            elif index == 6:
                video_id = shape_id(document, "media-video")
                audio_id = shape_id(document, "media-audio")
                bookmark_target_id = shape_id(document, "hero-媒体与书签")
                audio_picture = document.xpath(
                    ".//p:cNvPr[@name='media-audio']/ancestor::p:pic[1]",
                    namespaces=NS,
                )[0]
                audio_file = audio_picture.xpath(".//*[local-name()='videoFile']")[0]
                audio_file.tag = qname("a", "audioFile")
                audio_relation_id = audio_file.get(qname("r", "link"))
                rels_path = "ppt/slides/_rels/slide6.xml.rels"
                rels_document = etree.fromstring(archive.read(rels_path))
                audio_relation = rels_document.xpath(
                    ".//*[local-name()='Relationship'][@Id=$relation_id]",
                    relation_id=audio_relation_id,
                )[0]
                audio_relation.set(
                    "Type",
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio",
                )
                replacements[rels_path] = etree.tostring(
                    rels_document,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone=True,
                )
                media_extension = document.xpath(".//p14:media", namespaces=NS)[0]
                etree.SubElement(media_extension, qname("p14", "trim"), {"st": "100", "end": "1100"})
                bookmark_list = etree.SubElement(media_extension, qname("p14", "bmkLst"))
                etree.SubElement(bookmark_list, qname("p14", "bmk"), {"name": "middle", "time": "500"})
                document.insert(insert_at, media_timing(video_id, audio_id, bookmark_target_id))
            elif index == 7:
                document.insert(insert_at, repeat_timing(
                    shape_id(document, "hero-重复、往返与强调"),
                    shape_id(document, "trigger-重复、往返与强调"),
                    shape_id(document, "title-重复、往返与强调"),
                ))
            elif index == 8:
                document.insert(insert_at, transition("cut", 0, advance_ms=500))
                document.insert(insert_at + 1, indefinite_timing(shape_id(document, "hero-无限重复与停止")))
            elif index == 9:
                document.set("show", "0")
            elif index == 11:
                document.insert(insert_at, overlap_timing(shape_id(document, "hero-同属性冲突")))
            elif index == 12:
                document.insert(insert_at, paragraph_timing(shape_id(document, "title-文字逐段")))
            replacements[path] = etree.tostring(document, xml_declaration=True, encoding="UTF-8", standalone=True)
        target.parent.mkdir(parents=True, exist_ok=True)
        with ZipFile(target, "w", ZIP_DEFLATED) as output:
            for entry in archive.infolist():
                stable_entry = ZipInfo(entry.filename, date_time=(2026, 7, 11, 0, 0, 0))
                stable_entry.compress_type = ZIP_DEFLATED
                stable_entry.create_system = entry.create_system
                stable_entry.external_attr = entry.external_attr
                stable_entry.flag_bits = entry.flag_bits
                output.writestr(
                    stable_entry,
                    replacements.get(entry.filename, archive.read(entry.filename)),
                    compress_type=ZIP_DEFLATED,
                    compresslevel=9,
                )


def main() -> None:
    parser = argparse.ArgumentParser(description="生成 PPTX 播放回归素材。")
    parser.add_argument("--output", type=Path, default=OUTPUT)
    parser.add_argument(
        "--slides",
        default="1-12",
        help="只写入指定页面的测试结构，例如 1,2,6 或 1-6；默认全部。",
    )
    args = parser.parse_args()
    patched_slides: set[int] = set()
    for part in args.slides.split(","):
        start, separator, end = part.strip().partition("-")
        if not start:
            continue
        patched_slides.update(range(int(start), int(end) + 1) if separator else [int(start)])
    if not patched_slides or min(patched_slides) < 1 or max(patched_slides) > 12:
        parser.error("--slides 必须位于 1 到 12。")

    with TemporaryDirectory() as directory:
        base = Path(directory) / "base.pptx"
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        add_slide(prs, "单击动画", RGBColor(37, 99, 235))
        add_slide(prs, "淡化切换", RGBColor(5, 150, 105))
        add_slide(prs, "推进切换", RGBColor(217, 119, 6))
        add_slide(prs, "擦除与自动换页", RGBColor(220, 38, 38), "!!hero")
        add_slide(prs, "Morph 强身份", RGBColor(124, 58, 237), "!!hero")
        add_slide(prs, "媒体与书签", RGBColor(8, 145, 178))
        video = Path(directory) / "tiny.webm"
        audio = Path(directory) / "tiny.wav"
        poster = Path(directory) / "poster.png"
        image = numpy.zeros((180, 320, 3), dtype=numpy.uint8)
        image[:, :] = (8, 145, 178)
        with av.open(str(video), "w", options={"fflags": "+bitexact"}) as container:
            stream = container.add_stream("libvpx", rate=25)
            stream.codec_context.flags |= av.codec.context.Flags.bitexact
            stream.width = 320
            stream.height = 180
            stream.pix_fmt = "yuv420p"
            for _ in range(50):
                frame = av.VideoFrame.from_ndarray(image, format="rgb24")
                for packet in stream.encode(frame):
                    container.mux(packet)
            for packet in stream.encode():
                container.mux(packet)
        Image.fromarray(image).save(poster)
        audio_samples = (numpy.sin(2 * numpy.pi * 440 * numpy.arange(32000) / 16000) * 12000).astype(numpy.int16)
        with wave.open(str(audio), "wb") as audio_file:
            audio_file.setnchannels(1)
            audio_file.setsampwidth(2)
            audio_file.setframerate(16000)
            audio_file.writeframes(audio_samples.tobytes())
        movie = prs.slides[-1].shapes.add_movie(
            str(video), Inches(4.8), Inches(1.5), Inches(4), Inches(2.25),
            poster_frame_image=str(poster), mime_type="video/webm",
        )
        movie.name = "media-video"
        audio_shape = prs.slides[-1].shapes.add_movie(
            str(audio), Inches(8.9), Inches(4.5), Inches(0.5), Inches(0.5),
            poster_frame_image=str(poster), mime_type="audio/wav",
        )
        audio_shape.name = "media-audio"
        add_slide(prs, "重复、往返与强调", RGBColor(14, 116, 144))
        add_slide(prs, "无限重复与停止", RGBColor(190, 24, 93))
        add_slide(prs, "无限重复之后", RGBColor(71, 85, 105))
        add_slide(prs, "隐藏页之后", RGBColor(51, 65, 85))
        named_shape(prs, 9, "hero-隐藏页之后").click_action.target_slide = prs.slides[8]
        named_shape(prs, 9, "trigger-隐藏页之后").click_action.hyperlink.address = "https://example.com/pptx"
        add_slide(prs, "同属性冲突", RGBColor(126, 34, 206))
        add_slide(prs, "文字逐段", RGBColor(3, 105, 161))
        paragraph_box = named_shape(prs, 11, "title-文字逐段").text_frame
        paragraph_box.paragraphs[0].text = "第一段"
        paragraph_box.add_paragraph().text = "第二段"
        paragraph_box.add_paragraph().text = "第三段"
        prs.save(base)
        patch_archive(base, args.output, patched_slides)
    print(args.output)


if __name__ == "__main__":
    main()
